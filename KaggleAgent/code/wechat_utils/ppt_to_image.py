"""
Function: 替换PPT中的文本素材后，将PPT转换为图片
History:
20250616    HongfengAi
"""
import sys
from pathlib import Path
import argparse
from pptx import Presentation
import os
import win32com.client
import pythoncom

import sys, os
sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
from configs import OUTPUT_ROOT_PATH
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from utils import safe_title_func


class PPTToImage:
    def __init__(self):
        pass
    
    def split_run_at_placeholder(self, 
                                 paragraph, 
                                 run_index: int, 
                                 placeholder: str, 
                                 replacement: str):
        """在占位符处分割文本运行，保持格式"""
        run = paragraph.runs[run_index]
        text = run.text
        
        if f"[{placeholder}]" not in text:
            return False
        
        placeholder_full = f"[{placeholder}]"
        parts = text.split(placeholder_full)
        
        if len(parts) != 2:
            return False
        
        before_text, after_text = parts
        
        # 保存原始格式
        original_font = run.font
        
        # 更新当前run为前半部分
        run.text = before_text
        
        # 插入替换文本的新run
        replacement_run = paragraph.runs.add()
        replacement_run.text = replacement
        # 复制格式
        replacement_run.font.name = original_font.name
        replacement_run.font.size = original_font.size
        replacement_run.font.bold = original_font.bold
        replacement_run.font.italic = original_font.italic
        replacement_run.font.underline = original_font.underline
        
        # 插入后半部分的新run
        if after_text:
            after_run = paragraph.runs.add()
            after_run.text = after_text
            # 复制格式
            after_run.font.name = original_font.name
            after_run.font.size = original_font.size
            after_run.font.bold = original_font.bold
            after_run.font.italic = original_font.italic
            after_run.font.underline = original_font.underline
        
        return True
    
    def process_template(self, 
                        comp_type: str,
                        title: str, 
                        host: str, 
                        keywords: str):
        """高级替换方法，完美保持格式"""
        ppt_path = f"./code/wechat_utils/cover_template/{comp_type}_cover.pptx"
        prs = Presentation(ppt_path)
        
        replacements = {
            "title": title,
            "host": host,
            "keywords": keywords
        }
        
        # 遍历所有幻灯片和形状
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        # 从后往前处理runs，避免索引问题
                        run_indices = list(range(len(paragraph.runs)))
                        run_indices.reverse()
                        
                        for run_index in run_indices:
                            if run_index < len(paragraph.runs):
                                run = paragraph.runs[run_index]
                                
                                # 检查是否包含占位符
                                text = run.text
                                for placeholder, replacement in replacements.items():
                                    placeholder_full = f"[{placeholder}]"
                                    if placeholder_full in text:
                                        # 如果文本正好是占位符，直接替换
                                        if text.strip() == placeholder_full:
                                            run.text = replacement
                                        else:
                                            # 否则需要分割
                                            self.split_run_at_placeholder(
                                                paragraph, run_index, placeholder, replacement
                                            )
                                        break
        
        # 保存修改后的PPT
        safe_title = safe_title_func(title)
        output_path = f"{OUTPUT_ROOT_PATH}/{comp_type}/{safe_title}/cover"
        os.makedirs(output_path, exist_ok=True)
        prs.save(f"{output_path}/cover.pptx")


    def convert_ppt_to_image_with_style(self, comp_type: str, title: str) -> bool:
        """保持样式的PPT转图片"""

        safe_title = safe_title_func(title)
        ppt_path = f"{OUTPUT_ROOT_PATH}/{comp_type}/{safe_title}/cover/cover.pptx"
        image_path = f"{OUTPUT_ROOT_PATH}/{comp_type}/{safe_title}/cover/cover.png"
        print(f"正在转换: {Path(ppt_path).name} -> {Path(image_path).name}")
        
        # 检查输入文件是否存在
        if not Path(ppt_path).exists():
            print(f"✗ PPT文件不存在: {ppt_path}")
            return False
        
        # 创建输出目录
        Path(image_path).parent.mkdir(exist_ok=True)
        
        # 初始化COM
        pythoncom.CoInitialize()

        # 使用WPS转换
        wps = win32com.client.gencache.EnsureDispatch("Kwpp.Application")
        
        # 打开PPT文件
        ppt = wps.Presentations.Open(ppt_path)
        
        # 获取第一张幻灯片并导出
        first_slide = ppt.Slides.Item(1)
        first_slide.Export(image_path, "PNG")
        
        # 关闭文件和应用
        ppt.Close()
        wps.Quit()
        
        # 检查是否成功生成图片
        if Path(image_path).exists():
            print(f"✓ 成功生成: {image_path}")
            return True
        else:
            print("✗ 转换失败: 未能生成图片文件")
            return False     

        # 清理COM
        pythoncom.CoUninitialize()

    def main(self, comp_type: str, title: str, host: str, keywords: str):
        """替换PPT中的文本素材->将PPT转换为图片"""
        self.process_template(comp_type, title, host, keywords)
        self.convert_ppt_to_image_with_style(comp_type, title)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPT样式保持转换工具")
    parser.add_argument("--comp_type", type=str, default="comp_express", help="comp_express or comp_review")
    parser.add_argument("--title", type=str, default="CMI - Detect Behavior with Sensor Data", help="标题")
    parser.add_argument("--host", type=str, default="Child Mind Institute", help="主持人")
    parser.add_argument("--keywords", type=str, default="传感器数据处理、 行为识别、 机器学习、 时间序列分析", help="关键词")
    args = parser.parse_args()

    ppt_to_image = PPTToImage()
    ppt_to_image.main(args.comp_type, args.title, args.host, args.keywords)

