"""
Function: 替换PPT中的文本素材后，将PPT转换为图片
History:
20250616    HongfengAi
"""
import sys
from pathlib import Path
import argparse
from pptx import Presentation
import os
import win32com.client
import pythoncom

import sys, os
sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
from configs import OUTPUT_ROOT_PATH
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from utils import safe_title_func
from wechat_utils.upload_material import WeChatPermanentMaterialUploader

class PPTToImage:
    def __init__(self):
        self.wechat_uploader = WeChatPermanentMaterialUploader()
    
    def split_run_at_placeholder(self, 
                                 paragraph, 
                                 run_index: int, 
                                 placeholder: str, 
                                 replacement: str):
        """在占位符处分割文本运行，保持格式"""
        run = paragraph.runs[run_index]
        text = run.text
        
        if f"[{placeholder}]" not in text:
            return False
        
        placeholder_full = f"[{placeholder}]"
        parts = text.split(placeholder_full)
        
        if len(parts) != 2:
            return False
        
        before_text, after_text = parts
        
        # 保存原始格式
        original_font = run.font
        
        # 更新当前run为前半部分
        run.text = before_text
        
        # 插入替换文本的新run
        replacement_run = paragraph.runs.add()
        replacement_run.text = replacement
        # 复制格式
        replacement_run.font.name = original_font.name
        replacement_run.font.size = original_font.size
        replacement_run.font.bold = original_font.bold
        replacement_run.font.italic = original_font.italic
        replacement_run.font.underline = original_font.underline
        
        # 插入后半部分的新run
        if after_text:
            after_run = paragraph.runs.add()
            after_run.text = after_text
            # 复制格式
            after_run.font.name = original_font.name
            after_run.font.size = original_font.size
            after_run.font.bold = original_font.bold
            after_run.font.italic = original_font.italic
            after_run.font.underline = original_font.underline
        
        return True
    
    def process_template(self, 
                        year: str,
                        week: str):
        """高级替换方法，完美保持格式"""
        ppt_path = f"./code/wechat_utils/cover_template/app_express_cover.pptx"
        prs = Presentation(ppt_path)

        # week = week.replace("W", "")

        replacements = {
            "year": str(year),
            "week": str(week)
        }
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        # 检查段落是否有runs
                        if paragraph.runs:
                            # 将所有runs的文本拼接起来
                            paragraph_text = ""
                            for run in paragraph.runs:
                                paragraph_text += run.text
                            
                            # 在段落级别进行占位符替换
                            new_text = paragraph_text
                            text_changed = False
                            
                            for placeholder, replacement in replacements.items():
                                placeholder_full = f"[{placeholder}]"
                                if placeholder_full in new_text:    
                                    new_text = new_text.replace(placeholder_full, replacement)
                                    text_changed = True
                            
                            # 如果文本有变化，更新段落
                            if text_changed:
                                # 保存原始格式信息（从第一个非空run）
                                original_font = None
                                for run in paragraph.runs:
                                    if run.text.strip():
                                        original_font = run.font
                                        break
                                
                                # 清空所有runs
                                paragraph.clear()
                                
                                # 添加新的run with替换后的文本
                                new_run = paragraph.add_run()
                                new_run.text = new_text
                                
                                # 恢复原始格式
                                if original_font:
                                    try:
                                        # 复制基本字体属性
                                        if original_font.name:
                                            new_run.font.name = original_font.name
                                        if original_font.size:
                                            new_run.font.size = original_font.size
                                        if original_font.bold is not None:
                                            new_run.font.bold = original_font.bold
                                        if original_font.italic is not None:
                                            new_run.font.italic = original_font.italic
                                        
                                        # 复制字体颜色
                                        try:
                                            if hasattr(original_font.color, 'rgb') and original_font.color.rgb:
                                                new_run.font.color.rgb = original_font.color.rgb
                                            elif hasattr(original_font.color, 'theme_color') and original_font.color.theme_color is not None:
                                                new_run.font.color.theme_color = original_font.color.theme_color
                                        except:
                                            # 如果颜色复制失败，设置为黑色
                                            from pptx.dml.color import RGBColor
                                            new_run.font.color.rgb = RGBColor(0, 0, 0)
                                    except:
                                        pass  # 如果格式复制失败，忽略
                                
        
        # 保存修改后的PPT
        output_path = f"{OUTPUT_ROOT_PATH}/software_express/{year}_{week}/cover"
        os.makedirs(output_path, exist_ok=True)
        ppt_file_path = f"{output_path}/cover.pptx"
        prs.save(ppt_file_path)
        print(f"💾 PPT已保存到: {ppt_file_path}")
    


    def convert_ppt_to_image_with_style(self, year: str, week: str) -> bool:
        """保持样式的PPT转图片"""

        ppt_path = f"{OUTPUT_ROOT_PATH}/software_express/{year}_{week}/cover/cover.pptx"
        image_path = f"{OUTPUT_ROOT_PATH}/software_express/{year}_{week}/cover/cover.png"
        print(f"🔄 正在转换: {Path(ppt_path).name} -> {Path(image_path).name}")
        
        # 检查输入文件是否存在
        if not Path(ppt_path).exists():
            print(f"❌ PPT文件不存在: {ppt_path}")
            return False
        
        # 创建输出目录
        Path(image_path).parent.mkdir(exist_ok=True)
        
        # 初始化COM
        pythoncom.CoInitialize()

        # 使用WPS转换
        wps = win32com.client.gencache.EnsureDispatch("Kwpp.Application")
        
        # 打开PPT文件
        ppt = wps.Presentations.Open(ppt_path)
        
        # 获取第一张幻灯片并导出
        first_slide = ppt.Slides.Item(1)
        first_slide.Export(image_path, "PNG")
        
        # 关闭文件和应用
        ppt.Close()
        wps.Quit()

        
        # 检查是否成功生成图片
        if Path(image_path).exists():
            print(f"✓ 成功生成: {image_path}")
            self.wechat_uploader.upload_specific_file(image_path, title="app_express_cover")
            return True
        else:
            print("✗ 转换失败: 未能生成图片文件")
            return False

    def main(self, year: str, week: str):
        """替换PPT中的文本素材->将PPT转换为图片"""
        self.process_template(year, week)
        self.convert_ppt_to_image_with_style(year, week)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPT样式保持转换工具")
    parser.add_argument("--year", type=str, default="2025", help="年份")
    parser.add_argument("--week", type=str, default="26", help="周数")
    args = parser.parse_args()

    ppt_to_image = PPTToImage()
    ppt_to_image.main(args.year, args.week)

